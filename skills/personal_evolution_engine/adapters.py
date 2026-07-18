"""Deterministic, local-only adapters used by generated workflow recipes."""

from __future__ import annotations

import csv
import json
import shutil
from email.message import EmailMessage
from pathlib import Path
from typing import Any, Callable

_MAX_SOURCE_FILES = 200
_MAX_SOURCE_BYTES = 8_000_000


class WorkflowError(ValueError):
    """A safe, user-actionable workflow failure."""


def _safe_source_glob(workspace: Path, pattern: str) -> list[Path]:
    raw = str(pattern or "").strip()
    if not raw or Path(raw).is_absolute() or ".." in Path(raw).parts:
        raise WorkflowError("Один из источников указан небезопасно.")
    root = workspace.resolve()
    paths = sorted(path for path in root.glob(raw) if path.is_file())
    if not paths:
        raise WorkflowError("Не найдены данные, необходимые для выбранной автоматизации.")
    if len(paths) > _MAX_SOURCE_FILES:
        raise WorkflowError("Для одного запуска найдено слишком много исходных файлов.")
    total_bytes = 0
    for path in paths:
        resolved = path.resolve()
        if root not in resolved.parents:
            raise WorkflowError("Один из источников находится вне разрешённой папки.")
        total_bytes += path.stat().st_size
    if total_bytes > _MAX_SOURCE_BYTES:
        raise WorkflowError("Исходные данные слишком велики для одного тестового запуска.")
    return paths


def _safe_output(output_dir: Path, relative: str) -> Path:
    raw = str(relative or "").strip()
    path = Path(raw)
    if not raw or path.is_absolute() or ".." in path.parts:
        raise WorkflowError("Имя результата указано небезопасно.")
    root = output_dir.resolve()
    target = (root / path).resolve()
    if root != target and root not in target.parents:
        raise WorkflowError("Результат выходит за пределы локальной папки запуска.")
    target.parent.mkdir(parents=True, exist_ok=True)
    return target


def _records(context: dict[str, Any], key: str) -> list[dict[str, Any]]:
    value = context.get(str(key or ""))
    if not isinstance(value, list) or any(not isinstance(item, dict) for item in value):
        raise WorkflowError("Не удалось связать промежуточные данные выбранного процесса.")
    return [dict(item) for item in value]


def _setting_name(settings: dict[str, Any], key: str) -> str:
    value = str(settings.get(key) or "").strip()
    if not value:
        raise WorkflowError("В выбранной автоматизации не хватает настройки результата.")
    return value


def _coerce_scalar(value: Any) -> Any:
    if not isinstance(value, str):
        return value
    stripped = value.strip()
    if not stripped:
        return ""
    try:
        return int(stripped)
    except ValueError:
        try:
            return float(stripped.replace(",", "."))
        except ValueError:
            return stripped


def _load_json(context: dict[str, Any], settings: dict[str, Any], workspace: Path, output_dir: Path) -> list[Path]:
    del output_dir
    rows: list[dict[str, Any]] = []
    for path in _safe_source_glob(workspace, _setting_name(settings, "glob")):
        try:
            payload = json.loads(path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            raise WorkflowError(f"Не удалось прочитать синтетический источник «{path.name}».") from exc
        items = payload if isinstance(payload, list) else [payload]
        for item in items:
            if not isinstance(item, dict):
                raise WorkflowError(f"Источник «{path.name}» содержит неподдерживаемую запись.")
            rows.append({**item, "_source_file": path.name})
    context[_setting_name(settings, "save_as")] = rows
    return []


def _load_csv(context: dict[str, Any], settings: dict[str, Any], workspace: Path, output_dir: Path) -> list[Path]:
    del output_dir
    rows: list[dict[str, Any]] = []
    for path in _safe_source_glob(workspace, _setting_name(settings, "glob")):
        try:
            with path.open("r", encoding="utf-8-sig", newline="") as handle:
                for item in csv.DictReader(handle):
                    rows.append({key: _coerce_scalar(value) for key, value in item.items()})
        except (OSError, csv.Error) as exc:
            raise WorkflowError(f"Не удалось прочитать таблицу «{path.name}».") from exc
    context[_setting_name(settings, "save_as")] = rows
    return []


def _select_fields(
    context: dict[str, Any], settings: dict[str, Any], workspace: Path, output_dir: Path
) -> list[Path]:
    del workspace, output_dir
    fields = settings.get("fields")
    if not isinstance(fields, list) or not fields:
        raise WorkflowError("Не указаны поля для итоговой таблицы.")
    source = _records(context, _setting_name(settings, "from"))
    selected = [{str(field): row.get(str(field)) for field in fields} for row in source]
    context[_setting_name(settings, "save_as")] = selected
    return []


def _join_tables(context: dict[str, Any], settings: dict[str, Any], workspace: Path, output_dir: Path) -> list[Path]:
    del workspace, output_dir
    left = _records(context, _setting_name(settings, "left"))
    right = _records(context, _setting_name(settings, "right"))
    keys_raw = settings.get("on")
    keys = [str(item) for item in keys_raw] if isinstance(keys_raw, list) else [str(keys_raw or "")]
    if not all(keys):
        raise WorkflowError("Не указано поле для объединения данных.")
    suffix = str(settings.get("right_suffix") or "_right")
    index = {tuple(row.get(key) for key in keys): row for row in right}
    joined: list[dict[str, Any]] = []
    for left_row in left:
        output = dict(left_row)
        right_row = index.get(tuple(left_row.get(key) for key in keys), {})
        for key, value in right_row.items():
            if key in keys:
                continue
            output[key if key not in output else f"{key}{suffix}"] = value
        joined.append(output)
    context[_setting_name(settings, "save_as")] = joined
    return []


def _project_records(
    context: dict[str, Any], settings: dict[str, Any], workspace: Path, output_dir: Path
) -> list[Path]:
    del workspace, output_dir
    source = _records(context, _setting_name(settings, "from"))
    fields = settings.get("fields")
    if not isinstance(fields, dict) or not fields:
        raise WorkflowError("Не указано, какие поля подготовить для проверки.")
    result = []
    for row in source:
        projected: dict[str, Any] = {}
        for output_key, rule in fields.items():
            if isinstance(rule, str):
                projected[str(output_key)] = row.get(rule)
            elif isinstance(rule, dict) and "template" in rule:
                projected[str(output_key)] = _render(str(rule["template"]), row)
            elif isinstance(rule, dict) and "value" in rule:
                projected[str(output_key)] = rule["value"]
            else:
                raise WorkflowError("Одно из правил подготовки результата не поддерживается.")
        result.append(projected)
    context[_setting_name(settings, "save_as")] = result
    return []


def _compare_records(
    context: dict[str, Any], settings: dict[str, Any], workspace: Path, output_dir: Path
) -> list[Path]:
    del workspace, output_dir
    current = _records(context, _setting_name(settings, "current"))
    previous = _records(context, _setting_name(settings, "previous"))
    key = _setting_name(settings, "key")
    value = _setting_name(settings, "value")
    current_name = str(settings.get("current_name") or "current")
    previous_name = str(settings.get("previous_name") or "previous")
    delta_name = str(settings.get("delta_name") or "delta")
    previous_index = {row.get(key): row for row in previous}
    result = []
    for row in current:
        previous_row = previous_index.get(row.get(key), {})
        current_value = _coerce_scalar(row.get(value, 0))
        previous_value = _coerce_scalar(previous_row.get(value, 0))
        try:
            delta = round(float(current_value) - float(previous_value), 2)
        except (TypeError, ValueError):
            delta = ""
        result.append({
            key: row.get(key),
            current_name: current_value,
            previous_name: previous_value,
            delta_name: delta,
        })
    context[_setting_name(settings, "save_as")] = result
    return []


def _write_json(context: dict[str, Any], settings: dict[str, Any], workspace: Path, output_dir: Path) -> list[Path]:
    del workspace
    source = context.get(_setting_name(settings, "from"))
    target = _safe_output(output_dir, _setting_name(settings, "filename"))
    target.write_text(json.dumps(source, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return [target]


def _copy_files(context: dict[str, Any], settings: dict[str, Any], workspace: Path, output_dir: Path) -> list[Path]:
    del context
    directory = _safe_output(output_dir, _setting_name(settings, "directory"))
    directory.mkdir(parents=True, exist_ok=True)
    artifacts = []
    for source in _safe_source_glob(workspace, _setting_name(settings, "glob")):
        target = directory / source.name
        shutil.copy2(source, target)
        artifacts.append(target)
    return artifacts


def _column_specs(settings: dict[str, Any]) -> list[dict[str, Any]]:
    columns = settings.get("columns")
    if not isinstance(columns, list) or not columns:
        raise WorkflowError("Для итогового файла не указаны столбцы.")
    clean = []
    for column in columns:
        if not isinstance(column, dict) or not str(column.get("key") or "").strip():
            raise WorkflowError("Один из столбцов результата описан неполно.")
        clean.append({
            "key": str(column["key"]),
            "header": str(column.get("header") or column["key"]),
            "required": bool(column.get("required", False)),
            "number_format": str(column.get("number_format") or ""),
        })
    return clean


def _write_excel(context: dict[str, Any], settings: dict[str, Any], workspace: Path, output_dir: Path) -> list[Path]:
    del workspace
    try:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, Reference
        from openpyxl.chart.series import SeriesLabel
        from openpyxl.styles import Font, PatternFill
        from openpyxl.utils import get_column_letter
    except ImportError as exc:
        raise WorkflowError("Компонент создания Excel пока недоступен.") from exc
    rows = _records(context, _setting_name(settings, "from"))
    columns = _column_specs(settings)
    target = _safe_output(output_dir, _setting_name(settings, "filename"))
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = str(settings.get("sheet") or "Result")[:31]
    sheet.append([column["header"] for column in columns])
    for cell in sheet[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="C93545")
    missing_fill = PatternFill("solid", fgColor="FFF2CC")
    for row_index, row in enumerate(rows, 2):
        values = [row.get(column["key"]) for column in columns]
        sheet.append(values)
        for column_index, column in enumerate(columns, 1):
            value = values[column_index - 1]
            if column["number_format"]:
                sheet.cell(row=row_index, column=column_index).number_format = column["number_format"]
            if column["required"] and (value is None or str(value).strip() == ""):
                sheet.cell(row=row_index, column=column_index).fill = missing_fill
    for index, column in enumerate(columns, 1):
        max_length = max([len(str(column["header"]))] + [len(str(row.get(column["key"], ""))) for row in rows])
        sheet.column_dimensions[get_column_letter(index)].width = min(42, max(10, max_length + 2))
    sheet.freeze_panes = "A2"
    sheet.sheet_view.showGridLines = False
    chart = settings.get("chart")
    if isinstance(chart, dict) and rows:
        category_key = str(chart.get("category_key") or "")
        value_keys = [str(item) for item in chart.get("value_keys", [])]
        column_keys = [column["key"] for column in columns]
        if category_key in column_keys and value_keys and all(key in column_keys for key in value_keys):
            plot = BarChart()
            plot.title = str(chart.get("title") or "Summary")
            plot.y_axis.title = str(chart.get("value_title") or "Value")
            plot.x_axis.title = str(chart.get("category_title") or "Category")
            labels = chart.get("labels") if isinstance(chart.get("labels"), dict) else {}
            for key in value_keys:
                index = column_keys.index(key) + 1
                plot.add_data(Reference(sheet, min_col=index, min_row=1, max_row=len(rows) + 1), titles_from_data=True)
                column = columns[index - 1]
                plot.series[-1].tx = SeriesLabel(v=str(labels.get(key) or column["header"]))
            category_index = column_keys.index(category_key) + 1
            plot.set_categories(Reference(sheet, min_col=category_index, min_row=2, max_row=len(rows) + 1))
            plot.height = 8
            plot.width = 14
            sheet.add_chart(plot, str(chart.get("anchor") or "H2"))
    workbook.save(target)
    return [target]


def _write_presentation(
    context: dict[str, Any], settings: dict[str, Any], workspace: Path, output_dir: Path
) -> list[Path]:
    del workspace
    try:
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise WorkflowError("Компонент создания PowerPoint пока недоступен.") from exc
    rows = _records(context, _setting_name(settings, "from"))
    columns = _column_specs(settings)
    target = _safe_output(output_dir, _setting_name(settings, "filename"))
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = str(settings.get("title") or "Report")
    title_slide.placeholders[1].text = str(settings.get("subtitle") or "Prepared for review")

    table_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    table_slide.shapes.title.text = str(settings.get("table_title") or "Details")
    table_shape = table_slide.shapes.add_table(
        max(2, len(rows) + 1), len(columns), Inches(0.5), Inches(1.4), Inches(9), Inches(4.5)
    )
    table = table_shape.table
    for column_index, column in enumerate(columns):
        table.cell(0, column_index).text = column["header"]
    for row_index, row in enumerate(rows, 1):
        for column_index, column in enumerate(columns):
            table.cell(row_index, column_index).text = str(row.get(column["key"], ""))
    for cell in table.rows[0].cells:
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)

    chart = settings.get("chart")
    if isinstance(chart, dict) and rows:
        category_key = str(chart.get("category_key") or columns[0]["key"])
        value_keys = [str(item) for item in chart.get("value_keys", [])]
        if value_keys:
            chart_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            chart_slide.shapes.title.text = str(chart.get("title") or "Changes")
            chart_data = ChartData()
            chart_data.categories = [str(row.get(category_key, "")) for row in rows]
            for key in value_keys:
                values = []
                for row in rows:
                    try:
                        values.append(float(row.get(key, 0)))
                    except (TypeError, ValueError):
                        values.append(0.0)
                chart_data.add_series(str(chart.get("labels", {}).get(key, key)), values)
            chart_slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(0.8), Inches(1.5), Inches(8.5), Inches(4.8), chart_data,
            )
    presentation.save(target)
    return [target]


class _SafeValues(dict[str, Any]):
    def __missing__(self, key: str) -> str:
        return f"{{{key}}}"


def _render(template: str, row: dict[str, Any]) -> str:
    try:
        return template.format_map(_SafeValues(row))
    except (ValueError, AttributeError) as exc:
        raise WorkflowError("Шаблон результата содержит неподдерживаемую конструкцию.") from exc


def _safe_mail_filename(value: str, position: int) -> str:
    clean = "".join(character if character.isalnum() or character in "-_." else "_" for character in value)
    clean = clean.strip("._")[:70]
    return clean or f"draft_{position}"


def _draft_emails(context: dict[str, Any], settings: dict[str, Any], workspace: Path, output_dir: Path) -> list[Path]:
    del workspace
    rows = _records(context, _setting_name(settings, "from"))
    folder = _safe_output(output_dir, _setting_name(settings, "directory"))
    folder.mkdir(parents=True, exist_ok=True)
    to_field = _setting_name(settings, "to_field")
    name_field = str(settings.get("name_field") or to_field)
    subject_template = _setting_name(settings, "subject_template")
    body_template = _setting_name(settings, "body_template")
    artifacts = []
    for position, row in enumerate(rows, 1):
        recipient = str(row.get(to_field) or "").strip()
        if not recipient:
            raise WorkflowError("В одном из черновиков отсутствует получатель.")
        message = EmailMessage()
        message["From"] = str(settings.get("from_address") or "automation@example.test")
        message["To"] = recipient
        message["Subject"] = _render(subject_template, row)
        message.set_content(_render(body_template, row))
        filename = _safe_mail_filename(str(row.get(name_field) or recipient), position) + ".eml"
        target = folder / filename
        target.write_bytes(message.as_bytes())
        artifacts.append(target)
    return artifacts


_OPERATIONS: dict[str, Callable[[dict[str, Any], dict[str, Any], Path, Path], list[Path]]] = {
    "load_json": _load_json,
    "load_csv": _load_csv,
    "select_fields": _select_fields,
    "join_tables": _join_tables,
    "project_records": _project_records,
    "compare_records": _compare_records,
    "write_json": _write_json,
    "copy_files": _copy_files,
    "write_excel": _write_excel,
    "write_presentation": _write_presentation,
    "draft_emails": _draft_emails,
}


def supported_operations() -> frozenset[str]:
    return frozenset(_OPERATIONS)


def execute_workflow(steps: list[dict[str, Any]], workspace: Path, output_dir: Path) -> dict[str, Any]:
    if not isinstance(steps, list) or not steps or len(steps) > 30:
        raise WorkflowError("Автоматизация содержит неподдерживаемое число действий.")
    workspace = workspace.resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    context: dict[str, Any] = {}
    artifacts: list[Path] = []
    for step in steps:
        if not isinstance(step, dict):
            raise WorkflowError("Одно из действий автоматизации повреждено.")
        operation = str(step.get("operation") or "").strip()
        settings = step.get("settings")
        handler = _OPERATIONS.get(operation)
        if handler is None or not isinstance(settings, dict):
            raise WorkflowError("Автоматизация содержит действие, которое не прошло проверку.")
        artifacts.extend(handler(context, settings, workspace, output_dir))
    if not artifacts:
        raise WorkflowError("Автоматизация не создала результат для проверки.")
    return {"artifacts": [str(path) for path in artifacts], "context_keys": sorted(context)}
